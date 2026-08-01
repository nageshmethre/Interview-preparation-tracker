import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme Colors
    BG_COLOR = RGBColor(15, 23, 42)       # Dark Slate Navy
    CARD_BG = RGBColor(30, 41, 59)        # Lighter Card Slate
    ACCENT_INDIGO = RGBColor(99, 102, 241)# Electric Indigo
    ACCENT_CYAN = RGBColor(56, 189, 248)  # Cyan Highlight
    TEXT_WHITE = RGBColor(248, 250, 252)  # Bright White
    TEXT_MUTED = RGBColor(148, 163, 184) # Muted Gray

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.color.rgb = BG_COLOR
        return bg

    def add_header(slide, title_text, subtitle_text=""):
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        p.font.name = 'Arial'
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.size = Pt(14)
            p2.font.color.rgb = TEXT_MUTED
            p2.font.name = 'Arial'

    def add_card(slide, left, top, width, height, title, points):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = ACCENT_INDIGO
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        p.font.name = 'Arial'
        p.space_after = Pt(12)

        for pt in points:
            p_bullet = tf.add_paragraph()
            p_bullet.text = "•  " + pt
            p_bullet.font.size = Pt(14)
            p_bullet.font.color.rgb = TEXT_WHITE
            p_bullet.font.name = 'Arial'
            p_bullet.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    center_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
    center_card.fill.solid()
    center_card.fill.fore_color.rgb = CARD_BG
    center_card.line.color.rgb = ACCENT_INDIGO
    center_card.line.width = Pt(2)

    tf1 = center_card.text_frame
    tf1.word_wrap = True
    tf1.margin_top = Inches(0.6)

    p1 = tf1.paragraphs[0]
    p1.text = "🚀 PREPSPACE"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_CYAN
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf1.add_paragraph()
    p2.text = "Full-Stack AI-Powered Interview Preparation & Career Tracker"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)

    p3 = tf1.add_paragraph()
    p3.text = "Spring Boot 3.3.1  |  Java 21  |  PostgreSQL  |  Spring Security JWT  |  Vanilla JS SPA"
    p3.font.size = Pt(14)
    p3.font.color.rgb = ACCENT_INDIGO
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(20)

    p4 = tf1.add_paragraph()
    p4.text = "Presenter: Nagesh Methre  |  Production Live: api.stream-in.app"
    p4.font.size = Pt(14)
    p4.font.color.rgb = TEXT_MUTED
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(15)

    # -------------------------------------------------------------
    # SLIDE 2: Problem Statement & Industry Context
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "1. Problem Statement & Challenges", "The fragmented state of candidate interview preparation")

    add_card(slide2, 0.8, 1.8, 5.6, 5.0, "Current Industry Pain Points", [
        "Fragmented Tools: Candidates use 5+ apps (Notion, LeetCode, Excel, YouTube).",
        "No Unified Analytics: Difficult to identify weak topics across DSA & System Design.",
        "Inconsistent Revision: Lack of automated spaced repetition for flashcards.",
        "Manual Application Tracking: No structured pipeline for job applications."
    ])

    add_card(slide2, 6.8, 1.8, 5.6, 5.0, "Why Candidates Struggle", [
        "Zero Daily Consistency Feedback: Hard to maintain 30+ day coding streaks.",
        "Security Concerns: Insecure sessions and lack of multi-device controls.",
        "Lack of Mock Feedback: No instant scoring or response analysis.",
        "Goal: Need a single unified platform for tracking & practice."
    ])

    # -------------------------------------------------------------
    # SLIDE 3: Proposed Solution - PrepSpace
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "2. Proposed Solution — PrepSpace Portal", "Unified full-stack ecosystem for candidate success")

    add_card(slide3, 0.8, 1.8, 3.6, 5.0, "Centralized Dashboard", [
        "Real-time streak heatmap",
        "DSA & System Design status",
        "Daily question targets",
        "Enrolled courses & progress"
    ])

    add_card(slide3, 4.8, 1.8, 3.6, 5.0, "Practice & AI Engine", [
        "Spaced repetition flashcards",
        "AI Mock Interview Simulator",
        "Dynamic topic filter bank",
        "Automated readiness scores"
    ])

    add_card(slide3, 8.8, 1.8, 3.6, 5.0, "Enterprise Controls", [
        "Full 15-module settings panel",
        "Active device session audit",
        "JWT + BCrypt security",
        "PDF & Excel exports"
    ])

    # -------------------------------------------------------------
    # SLIDE 4: Architecture & Tech Stack
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "3. System Architecture & Tech Stack", "Modern high-performance enterprise micro-monolith")

    add_card(slide4, 0.8, 1.8, 3.6, 5.0, "Frontend Tier", [
        "Vanilla HTML5 / Custom CSS3",
        "Glassmorphism Dark Mode",
        "Modular SPA Router (app.js)",
        "iText & POI Client Exports"
    ])

    add_card(slide4, 4.8, 1.8, 3.6, 5.0, "Backend Tier", [
        "Spring Boot 3.3.1 (Java 21)",
        "Spring Security (Stateless JWT)",
        "Spring Data JPA & Hibernate",
        "Apache Tomcat on Port 8080"
    ])

    add_card(slide4, 8.8, 1.8, 3.6, 5.0, "Database & Cloud", [
        "PostgreSQL Relational DB",
        "Render Cloud Container Host",
        "HikariCP Connection Pool",
        "SSL Mode Encrypted Traffic"
    ])

    # -------------------------------------------------------------
    # SLIDE 5: Security Architecture
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "4. Security & Authentication Blueprint", "Multi-layered protection for user credentials and sessions")

    add_card(slide5, 0.8, 1.8, 5.6, 5.0, "Authentication & Tokens", [
        "Stateless JWT Tokens: Standard HMAC-SHA256 signed tokens.",
        "BCrypt Password Hashing: Strong salted hashing via BCryptPasswordEncoder.",
        "Role-Based Access Control: Security rules (@PreAuthorize ROLE_STUDENT).",
        "CORS Security: Strict origin verification for production."
    ])

    add_card(slide5, 6.8, 1.8, 5.6, 5.0, "Session & Lockout Shield", [
        "Account Lockout Protection: 5 failed attempts triggers 15-min lock.",
        "Device Session Control: Track IP, Browser User-Agent, and active sessions.",
        "Remote Revocation: Revoke session tokens remotely from Settings.",
        "Password Rules: Enforce 6+ characters minimum length."
    ])

    # -------------------------------------------------------------
    # SLIDE 6: 15-Module Settings Suite
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "5. 15-Module Enterprise Settings Suite", "Comprehensive control center with 31 persisted attributes")

    add_card(slide6, 0.8, 1.8, 3.6, 5.0, "Identity & Security", [
        "1. Profile Settings",
        "2. Security & 2FA Settings",
        "3. Connected Accounts",
        "4. Devices & Sessions Audit",
        "5. Data & Privacy Controls"
    ])

    add_card(slide6, 4.8, 1.8, 3.6, 5.0, "Preferences & UI", [
        "6. Appearance (Themes/Fonts)",
        "7. Notification Preferences",
        "8. Language & Region Controls",
        "9. Learning Goals & Tech Stack",
        "10. Career & Salary Targets"
    ])

    add_card(slide6, 8.8, 1.8, 3.6, 5.0, "Advanced & System", [
        "11. Dashboard Widget Toggles",
        "12. Subscription Tier Status",
        "13. Data Import / Export",
        "14. Developer API Keys & Webhooks",
        "15. System About & Engine Status"
    ])

    # -------------------------------------------------------------
    # SLIDE 7: Database & Schema Design
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "6. Database Schema & Data Persistence", "PostgreSQL relational schema with Hibernate auto-migrations")

    add_card(slide7, 0.8, 1.8, 5.6, 5.0, "Core Entities & Mappings", [
        "users Table: Primary user credentials, roles, and timestamps.",
        "user_settings Table: 31 columns storing layout, themes, goals, & API keys.",
        "job_applications Table: Applications pipeline, company, role, salary.",
        "mock_interviews Table: Session scores, feedback, dynamic questions."
    ])

    add_card(slide7, 6.8, 1.8, 5.6, 5.0, "Relational Integrity", [
        "Foreign Key Constraints: ON DELETE CASCADE for clean cleanup.",
        "Unique Email Constraints: Guaranteed single identity per email.",
        "HikariCP Connection Pool: Optimized high-concurrency connections.",
        "Hibernate DDL Auto: Dynamic schema synchronization."
    ])

    # -------------------------------------------------------------
    # SLIDE 8: AI Interview & Flashcard Engine
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "7. AI Interviewer & Flashcard Engine", "Smart algorithms driving retention and interview readiness")

    add_card(slide8, 0.8, 1.8, 5.6, 5.0, "AI Mock Interview Simulator", [
        "Adaptive Prompting: Generates DSA, Java, & System Design questions.",
        "Instant Scoring: Evaluates answer accuracy, clarity, and complexity.",
        "Feedback Reports: Provides strengths, weaknesses, and improvement steps."
    ])

    add_card(slide8, 6.8, 1.8, 5.6, 5.0, "Spaced-Repetition Flashcards", [
        "SuperMemo SM-2 Principle: Calculates review intervals by score.",
        "Dynamic Due Dates: Schedules cards for 1, 3, 7, or 14-day intervals.",
        "CS Fundamentals: Covers algorithms, concurrency, and OOP design."
    ])

    # -------------------------------------------------------------
    # SLIDE 9: Analytics & Export Engine
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_header(slide9, "8. Analytics & Report Generation", "Data-driven insights and multi-format document exports")

    add_card(slide9, 0.8, 1.8, 5.6, 5.0, "Visual Performance Analytics", [
        "GitHub-Style Activity Heatmap: Visual tracking of daily practice.",
        "Daily Streak Counter: Encourages 30+ day coding habits.",
        "Topic Mastery Gauges: Identifies weak topics for targeted revision."
    ])

    add_card(slide9, 6.8, 1.8, 5.6, 5.0, "Document Export Capabilities", [
        "iText PDF Generator: Downloads formatted readiness reports.",
        "Apache POI Excel Export: Exports full job application logs.",
        "JSON/CSV Data Backups: Portable user data exports."
    ])

    # -------------------------------------------------------------
    # SLIDE 10: Cloud Deployment & CI/CD
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)
    add_header(slide10, "9. Cloud Infrastructure & Deployment", "Production live hosting on Render Cloud container environment")

    add_card(slide10, 0.8, 1.8, 5.6, 5.0, "Render Cloud Container", [
        "Dockerized Maven Build: Multi-stage openjdk:21-slim container.",
        "Automatic Auto-Deploy: Triggers builds directly from main branch.",
        "Production Domain: Serves live API on https://api.stream-in.app."
    ])

    add_card(slide10, 6.8, 1.8, 5.6, 5.0, "Managed Cloud Database", [
        "PostgreSQL Cloud Instance: SSL-encrypted connection URI.",
        "Dynamic Env Vars: Standard SPRING_DATASOURCE_DRIVER_CLASS_NAME.",
        "Zero Downtime: Automatic health checks and restarts."
    ])

    # -------------------------------------------------------------
    # SLIDE 11: Live Demonstration Walkthrough
    # -------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11)
    add_header(slide11, "10. Live Demonstration Flow", "Step-by-step user journey across the application")

    add_card(slide11, 0.8, 1.8, 3.6, 5.0, "Step 1: Auth & Signup", [
        "User registers account",
        "JWT token issued on login",
        "Session stored in SPA state"
    ])

    add_card(slide11, 4.8, 1.8, 3.6, 5.0, "Step 2: Practice & AI", [
        "Explore questions bank",
        "Solve DSA & rate flashcards",
        "Take AI Mock Interview"
    ])

    add_card(slide11, 8.8, 1.8, 3.6, 5.0, "Step 3: Settings & Export", [
        "Open 15-module Settings",
        "Modify location & AI model",
        "Export readiness PDF report"
    ])

    # -------------------------------------------------------------
    # SLIDE 12: Future Roadmap & Vision
    # -------------------------------------------------------------
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12)
    add_header(slide12, "11. Future Expansion Roadmap", "Scalability horizons for upcoming platform releases")

    add_card(slide12, 0.8, 1.8, 3.6, 5.0, "Phase 1: Code Sandbox", [
        "Monaco Editor integration",
        "Multi-language compiler",
        "C++, Java, Python support",
        "Automated test suite"
    ])

    add_card(slide12, 4.8, 1.8, 3.6, 5.0, "Phase 2: Voice AI", [
        "Real-time audio interviews",
        "WebRTC speech stream",
        "Voice sentiment analysis",
        "Live technical coaching"
    ])

    add_card(slide12, 8.8, 1.8, 3.6, 5.0, "Phase 3: Peer Mocks", [
        "P2P Candidate Matching",
        "Collaborative Whiteboard",
        "Live shared coding buffer",
        "Peer feedback rubrics"
    ])

    # -------------------------------------------------------------
    # SLIDE 13: Conclusion & Q&A
    # -------------------------------------------------------------
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide13)
    add_header(slide13, "12. Conclusion & Summary", "PrepSpace — Empowering candidate career transformations")

    summary_card = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.8), Inches(10.333), Inches(4.8))
    summary_card.fill.solid()
    summary_card.fill.fore_color.rgb = CARD_BG
    summary_card.line.color.rgb = ACCENT_INDIGO
    summary_card.line.width = Pt(2)

    tf_s = summary_card.text_frame
    tf_s.word_wrap = True
    tf_s.margin_top = Inches(0.4)
    tf_s.margin_left = Inches(0.4)

    ps1 = tf_s.paragraphs[0]
    ps1.text = "🎯 Key Takeaways"
    ps1.font.size = Pt(22)
    ps1.font.bold = True
    ps1.font.color.rgb = ACCENT_CYAN
    ps1.space_after = Pt(12)

    bullets = [
        "Unified Full-Stack Solution: Solves fragmented interview prep with a cohesive Spring Boot + Vanilla JS architecture.",
        "Enterprise Security & Controls: Robust JWT + BCrypt auth with 15 granular settings modules.",
        "Production Tested & Live: Fully deployed on Render Cloud + PostgreSQL DB serving real API requests.",
        "Extensible Architecture: Built to easily integrate code sandboxes, WebRTC voice interviews, and P2P matching."
    ]
    for b in bullets:
        pb = tf_s.add_paragraph()
        pb.text = "•  " + b
        pb.font.size = Pt(16)
        pb.font.color.rgb = TEXT_WHITE
        pb.space_after = Pt(10)

    p_end = tf_s.add_paragraph()
    p_end.text = "Thank you! Any Questions?"
    p_end.font.size = Pt(20)
    p_end.font.bold = True
    p_end.font.color.rgb = ACCENT_CYAN
    p_end.alignment = PP_ALIGN.CENTER
    p_end.space_before = Pt(15)

    output_filename = "Interview_Preparation_Tracker_Presentation.pptx"
    prs.save(output_filename)

    artifact_dir = r"C:\Users\Nagesh\.gemini\antigravity\brain\69d1292e-0547-405f-9445-386116a1098c"
    if os.path.exists(artifact_dir):
        artifact_path = os.path.join(artifact_dir, "Interview_Preparation_Tracker_Presentation.pptx")
        prs.save(artifact_path)
        print(f"Presentation saved successfully to {output_filename} and {artifact_path}")
    else:
        print(f"Presentation saved successfully to {output_filename}")

if __name__ == "__main__":
    create_presentation()
